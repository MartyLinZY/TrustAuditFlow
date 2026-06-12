#!/usr/bin/env python3
from __future__ import annotations

import argparse
import json
import statistics
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]

BG = RGBColor(10, 18, 36)
BG2 = RGBColor(16, 27, 52)
PANEL = RGBColor(24, 38, 70)
CYAN = RGBColor(54, 211, 209)
BLUE = RGBColor(92, 142, 255)
PURPLE = RGBColor(151, 102, 255)
GREEN = RGBColor(88, 211, 132)
WHITE = RGBColor(244, 248, 255)
MUTED = RGBColor(178, 190, 214)
WARN = RGBColor(255, 185, 92)
RED = RGBColor(255, 100, 115)
FONT = "PingFang SC"


def load_json(path: Path) -> dict[str, Any]:
    return json.loads(path.read_text(encoding="utf-8"))


def cm(v: float) -> float:
    return v / 2.54


def set_fill(shape, color: RGBColor, transparency: float = 0.0) -> None:
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    if transparency:
        shape.fill.transparency = transparency
    shape.line.color.rgb = RGBColor(53, 73, 118)


def add_bg(slide, title: str, subtitle: str | None = None) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(0.16))
    bar.fill.solid()
    bar.fill.fore_color.rgb = CYAN
    bar.line.fill.background()
    tx = slide.shapes.add_textbox(Inches(0.55), Inches(0.35), Inches(8.8), Inches(0.55))
    p = tx.text_frame.paragraphs[0]
    p.text = title
    p.font.name = FONT
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = WHITE
    if subtitle:
        sub = slide.shapes.add_textbox(Inches(0.57), Inches(0.9), Inches(8.5), Inches(0.35))
        sp = sub.text_frame.paragraphs[0]
        sp.text = subtitle
        sp.font.name = FONT
        sp.font.size = Pt(10)
        sp.font.color.rgb = MUTED


def add_footer(slide, page: int) -> None:
    footer = slide.shapes.add_textbox(Inches(0.55), Inches(7.06), Inches(12.2), Inches(0.22))
    p = footer.text_frame.paragraphs[0]
    p.text = f"TrustFlowAudit | 数据流通场景下基于区块链的分布式校验与审计 | {page:02d}"
    p.font.name = FONT
    p.font.size = Pt(8)
    p.font.color.rgb = RGBColor(120, 136, 165)
    p.alignment = PP_ALIGN.RIGHT


def add_title(slide, title: str, subtitle: str, meta: str) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG
    slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.0), Inches(0.0), Inches(13.333), Inches(0.18)).fill.solid()
    shape = slide.shapes[-1]
    shape.fill.fore_color.rgb = CYAN
    shape.line.fill.background()
    left = slide.shapes.add_textbox(Inches(0.7), Inches(1.45), Inches(10.8), Inches(1.4))
    tf = left.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.name = FONT
    p.font.size = Pt(34)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p2 = tf.add_paragraph()
    p2.text = subtitle
    p2.font.name = FONT
    p2.font.size = Pt(17)
    p2.font.color.rgb = CYAN
    p2.space_before = Pt(12)
    m = slide.shapes.add_textbox(Inches(0.75), Inches(5.85), Inches(7.5), Inches(0.5))
    mp = m.text_frame.paragraphs[0]
    mp.text = meta
    mp.font.name = FONT
    mp.font.size = Pt(12)
    mp.font.color.rgb = MUTED
    for i, color in enumerate([CYAN, BLUE, PURPLE]):
        circ = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(9.2 + i * 0.65), Inches(4.2 - i * 0.45), Inches(1.7), Inches(1.7))
        circ.fill.solid()
        circ.fill.fore_color.rgb = color
        circ.fill.transparency = 45
        circ.line.fill.background()


def text_box(slide, x, y, w, h, text: str, size: int = 13, color: RGBColor = WHITE, bold: bool = False) -> None:
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.02)
    tf.margin_right = Inches(0.02)
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = FONT
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color


def card(slide, x, y, w, h, title: str, body: str, accent: RGBColor = CYAN) -> None:
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    set_fill(shape, PANEL)
    shape.adjustments[0] = 0.08
    stripe = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(0.08), Inches(h))
    stripe.fill.solid()
    stripe.fill.fore_color.rgb = accent
    stripe.line.fill.background()
    text_box(slide, x + 0.22, y + 0.18, w - 0.4, 0.35, title, 14, WHITE, True)
    text_box(slide, x + 0.22, y + 0.62, w - 0.4, h - 0.75, body, 10, MUTED)


def metric(slide, x, y, label: str, value: str, note: str, color: RGBColor = CYAN) -> None:
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(2.7), Inches(1.25))
    set_fill(shape, PANEL)
    text_box(slide, x + 0.18, y + 0.15, 2.25, 0.22, label, 9, MUTED)
    text_box(slide, x + 0.18, y + 0.42, 2.25, 0.4, value, 19, color, True)
    text_box(slide, x + 0.18, y + 0.88, 2.25, 0.25, note, 8, MUTED)


def add_table(slide, x, y, w, h, headers: list[str], rows: list[list[str]], widths: list[float] | None = None) -> None:
    table = slide.shapes.add_table(len(rows) + 1, len(headers), Inches(x), Inches(y), Inches(w), Inches(h)).table
    if widths:
        for i, width in enumerate(widths):
            table.columns[i].width = Inches(width)
    for c, header in enumerate(headers):
        cell = table.cell(0, c)
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(30, 58, 100)
        p = cell.text_frame.paragraphs[0]
        p.text = header
        p.font.name = FONT
        p.font.size = Pt(9)
        p.font.bold = True
        p.font.color.rgb = WHITE
    for r, row in enumerate(rows, start=1):
        for c, value in enumerate(row):
            cell = table.cell(r, c)
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(18, 31, 59) if r % 2 else RGBColor(22, 37, 70)
            p = cell.text_frame.paragraphs[0]
            p.text = value
            p.font.name = FONT
            p.font.size = Pt(8)
            p.font.color.rgb = WHITE if c == 0 else MUTED
            p.alignment = PP_ALIGN.CENTER if c else PP_ALIGN.LEFT


def flow(slide, labels: list[str], x: float, y: float, w: float = 1.7, gap: float = 0.32) -> None:
    colors = [CYAN, BLUE, PURPLE, GREEN, WARN]
    for i, label in enumerate(labels):
        sx = x + i * (w + gap)
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(sx), Inches(y), Inches(w), Inches(0.82))
        set_fill(box, RGBColor(22, 42, 78))
        box.line.color.rgb = colors[i % len(colors)]
        box.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = box.text_frame.paragraphs[0]
        p.text = label
        p.font.name = FONT
        p.font.size = Pt(9)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER
        if i < len(labels) - 1:
            line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(sx + w), Inches(y + 0.41), Inches(sx + w + gap), Inches(y + 0.41))
            line.line.color.rgb = RGBColor(88, 111, 154)
            line.line.width = Pt(1.2)


def chart_bar(slide, x, y, w, h, categories: list[str], series: list[tuple[str, list[float]]], title: str = "") -> None:
    data = CategoryChartData()
    data.categories = categories
    for name, values in series:
        data.add_series(name, values)
    chart = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(x), Inches(y), Inches(w), Inches(h), data).chart
    chart.has_title = bool(title)
    if title:
        chart.chart_title.text_frame.text = title
    chart.has_legend = len(series) > 1
    if chart.has_legend:
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False
    chart.value_axis.tick_labels.font.size = Pt(8)
    chart.category_axis.tick_labels.font.size = Pt(8)


def chart_line(slide, x, y, w, h, labels: list[str], values: list[float], title: str = "") -> None:
    data = CategoryChartData()
    data.categories = labels
    data.add_series("Accuracy", values)
    chart = slide.shapes.add_chart(XL_CHART_TYPE.LINE_MARKERS, Inches(x), Inches(y), Inches(w), Inches(h), data).chart
    chart.has_title = bool(title)
    if title:
        chart.chart_title.text_frame.text = title
    chart.has_legend = False
    chart.value_axis.tick_labels.font.size = Pt(8)
    chart.category_axis.tick_labels.font.size = Pt(8)


def notes_section(page: int, title: str, bullets: list[str]) -> str:
    lines = [f"## 第 {page} 页：{title}", ""]
    lines.extend([f"- {item}" for item in bullets])
    lines.append("")
    return "\n".join(lines)


def build_deck(output_dir: Path) -> tuple[Path, Path]:
    default = load_json(ROOT / "container_results_mnist" / "consortium_results.json")
    low = load_json(ROOT / "container_results_mnist_lowtrust" / "consortium_results.json")
    fault = load_json(ROOT / "container_results_mnist_fault" / "consortium_results.json")
    sim = load_json(ROOT / "results" / "results.json")
    scenarios = [("默认", default), ("低可信", low), ("故障", fault)]
    lat_avg = [round(statistics.mean([row["latency_ms"] for row in item["blocks"]]), 2) for _, item in scenarios]
    lat_max = [round(max(row["latency_ms"] for row in item["blocks"]), 2) for _, item in scenarios]
    acc_curve = [row["accuracy"] for row in default["metrics"]]
    acc_labels = [str(row["round"]) for row in default["metrics"]]
    summary = sim["strategy_summary"]
    baseline = next(item for item in summary if item["strategy"] == "blop_full_pbft")
    ours = next(item for item in summary if item["strategy"] == "trustflowaudit")
    record_reduction = 1.0 - ours["total_chain_records"] / baseline["total_chain_records"]
    latency_reduction = 1.0 - ours["avg_amortized_event_commit_ms"] / baseline["avg_amortized_event_commit_ms"]
    storage_reduction = 1.0 - ours["evidence_storage_overhead"] / baseline["evidence_storage_overhead"]

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    notes: list[str] = []

    slide = prs.slides.add_slide(blank)
    add_title(
        slide,
        "TrustFlowAudit：面向数据流通的区块链分布式校验与审计原型",
        "BLoP 可实现机制复现、ZooKeeper 联盟链验证与真实 MNIST 初步实验",
        "20 分钟阶段汇报 | 研究背景 · 方法设计 · 实验设计 · 初步结果",
    )
    notes.append(notes_section(1, "封面", [
        "这一页先给出汇报主题：工作不是做一个完整生产系统，而是围绕数据流通中的可信校验与审计，复现 BLoP 中可实现的机制，并进一步形成 TrustFlowAudit 原型。",
        "本次新增的关键进展是真实 MNIST 数据已经接入，并且在 ZooKeeper + Docker Compose 的容器化联盟链路径上完成了三组实验验证。",
        "汇报结构会按背景、动机、方法、实验设计、初步结果和后续计划展开。",
    ]))

    slide = prs.slides.add_slide(blank)
    add_bg(slide, "汇报主线", "从问题定义到可验证原型")
    flow(slide, ["数据流通可信问题", "传统审计能力边界", "BLoP机制复现", "TrustFlowAudit方法", "真实MNIST验证"], 0.75, 2.0, 2.05, 0.32)
    card(slide, 0.8, 3.45, 3.6, 1.55, "核心问题", "数据跨主体流通后，结果是否被正确计算、证据是否可追溯、异常节点是否能被排除。", CYAN)
    card(slide, 4.85, 3.45, 3.6, 1.55, "技术路径", "链下计算和证据摘要，链上记录关键哈希、投票与处置状态，降低全量上链成本。", BLUE)
    card(slide, 8.9, 3.45, 3.6, 1.55, "本次进展", "使用真实 MNIST 数据完成容器化联盟链实验，覆盖默认、低可信节点和故障节点场景。", GREEN)
    add_footer(slide, 2)
    notes.append(notes_section(2, "汇报主线", [
        "这页说明汇报的逻辑链条：先解释为什么数据流通需要可信审计，再说明传统 PDP/PoR 和普通区块链存证各自解决什么、不能解决什么。",
        "然后介绍 BLoP 机制中可以工程化复现的部分，包括可信评分、异常记录、PBFT 替代 PoW、链上防篡改等。",
        "最后落到 TrustFlowAudit：用轻量委员会、哈希摘要、Merkle root、纠删码机制模型和 ZooKeeper 容器联盟链做一个可验证闭环。",
    ]))

    slide = prs.slides.add_slide(blank)
    add_bg(slide, "研究背景：数据流通需要可验证的过程信任", "数据从“持有”走向“流通”和“协同计算”后，单点日志难以支撑责任闭环")
    card(slide, 0.75, 1.35, 3.7, 1.6, "数据不出域", "联邦学习、隐私计算和数据空间弱化了原始数据直接交换，但也让审计对象从文件完整性扩展到计算过程和结果摘要。", CYAN)
    card(slide, 4.8, 1.35, 3.7, 1.6, "跨主体协作", "多机构参与时，单方日志难以让其他参与方信任，尤其需要可共同验证的时间戳、哈希和责任记录。", BLUE)
    card(slide, 8.85, 1.35, 3.7, 1.6, "异常可追溯", "数据错误、模型异常、节点失效和低可信节点都需要被记录、定位和处置，而不是只给最终结果。", PURPLE)
    text_box(slide, 1.0, 4.0, 11.4, 0.75, "因此，数据流通场景中的审计问题不是单纯“文件有没有被改”，而是“谁在什么状态下对什么证据做了什么操作，其他节点如何验证并形成可追溯共识”。", 18, WHITE, True)
    add_footer(slide, 3)
    notes.append(notes_section(3, "研究背景", [
        "数据流通场景下，原始数据通常不希望直接集中到一处，而是通过隐私计算、联邦学习或者安全多方计算完成协同。",
        "这会带来一个新问题：我们不仅要知道数据本身有没有被改，还要知道计算过程是否可信、结果摘要是否可验证、异常责任是否能追溯。",
        "所以本课题关注的是过程级审计，而不是单纯的存储完整性验证。",
    ]))

    slide = prs.slides.add_slide(blank)
    add_bg(slide, "研究动机：传统审计与区块链能力需要结合", "PDP/PoR 解决远程数据完整性，区块链补充多方共识和不可抵赖记录")
    add_table(slide, 0.7, 1.3, 12.0, 3.0,
              ["路线", "主要解决", "不足", "本工作借鉴"],
              [
                  ["PDP/PoR", "远程数据完整性、抽样验证", "不直接处理多主体责任和链上治理", "抽样校验、哈希承诺"],
                  ["区块链存证", "不可篡改记录、多方共识", "全量上链成本高，计算过程仍需链下完成", "摘要存证、PBFT确认"],
                  ["BLoP", "可信节点、白名单、PBFT和治理闭环", "依赖可信计算设定，生产系统复刻成本高", "复现可实现机制并轻量化"],
                  ["TrustFlowAudit", "轻量委员会、分级存证、纠删码审计", "当前仍是原型和机制验证", "形成短周期可验证闭环"],
              ],
              widths=[1.7, 2.8, 3.7, 3.8])
    text_box(slide, 1.0, 5.1, 11.2, 0.55, "目标：不追求一次性实现完整可信计算系统，而是先验证“低成本审计委员会 + 证据摘要 + 链上确认”的可行性。", 17, CYAN, True)
    add_footer(slide, 4)
    notes.append(notes_section(4, "研究动机", [
        "传统 PDP/PoR 很适合做远程数据完整性验证，但它主要解决数据是否被保存和可恢复，不天然处理多方协作中的节点治理。",
        "区块链提供不可篡改和多方确认，但如果所有证据都上链，成本和延迟会很高。",
        "BLoP 给了一个方向：用可信节点集合和 PBFT 替代 PoW，但我们当前不做 TPM/TEE，而是把可实现机制抽象出来，形成轻量实验闭环。",
    ]))

    slide = prs.slides.add_slide(blank)
    add_bg(slide, "BLoP 可实现机制与本原型取舍", "明确不使用 TPM/TEE，先验证软件可信评分和审计流程是否能闭环")
    flow(slide, ["可信评分", "白名单/异常", "哈希链/Merkle", "PBFT确认", "处置记录"], 0.95, 1.55, 1.85, 0.35)
    metric(slide, 0.9, 3.15, "可信异常事件", f"{sim['totals']['trusted_anomaly_events']:,}", "按 BLoP 统计校准", CYAN)
    metric(slide, 3.85, 3.15, "防篡改事件", f"{sim['totals']['tamper_proof_events']:,}", "进入哈希链与 Merkle", BLUE)
    metric(slide, 6.8, 3.15, "白名单阻断", f"{sim['governance_summary']['total_whitelist_blocks']:,}", "治理事件模拟", GREEN)
    metric(slide, 9.75, 3.15, "链上处置记录", f"{sim['governance_summary']['total_chain_policy_records']:,}", "批量上链模型", PURPLE)
    card(slide, 1.0, 5.0, 5.3, 1.0, "保留", "可信评分、异常事件、PBFT确认、哈希链、Merkle root、节点治理。", CYAN)
    card(slide, 7.0, 5.0, 5.3, 1.0, "替代", "不使用 TPM/TEE/SGX，不做真实 PCR 远程证明，用软件评分和可审计日志替代。", WARN)
    add_footer(slide, 5)
    notes.append(notes_section(5, "BLoP 机制取舍", [
        "这页说明本原型和 BLoP 的关系：不是完整复刻生产环境，而是复现其中可以短期实现和验证的机制。",
        "保留的机制包括可信评分、异常检测、防篡改日志、PBFT确认和治理记录。",
        "明确替代的部分是 TPM/TEE/SGX 等可信硬件，因为当前目标是低成本、短周期、小论文级验证。",
    ]))

    slide = prs.slides.add_slide(blank)
    add_bg(slide, "方法：TrustFlowAudit 轻量审计闭环", "可信评分驱动审计委员会，链下保留证据，链上固化摘要和投票结果")
    flow(slide, ["数据/模型更新", "证据哈希", "可信委员会", "PBFT投票", "链上区块", "异常处置"], 0.45, 2.0, 1.75, 0.28)
    card(slide, 0.9, 3.6, 3.6, 1.35, "链下", "保存原始证据、模型更新、分片信息和计算过程，只把摘要和关键状态上链。", BLUE)
    card(slide, 4.85, 3.6, 3.6, 1.35, "链上", "记录 payload hash、Merkle root、prev hash、投票节点、quorum 与最终区块哈希。", CYAN)
    card(slide, 8.8, 3.6, 3.6, 1.35, "治理", "低可信节点排除；节点故障时保留剩余可信节点投票，满足 quorum 后继续提交。", GREEN)
    add_footer(slide, 6)
    notes.append(notes_section(6, "方法总览", [
        "TrustFlowAudit 的方法可以概括为三句话：链下保留大证据，链上保存可验证摘要；用可信评分选择委员会；通过 PBFT 风格投票形成提交结果。",
        "每个区块里记录 payload hash、Merkle root、前一区块哈希和投票结果，因此可以检查哈希链是否连续。",
        "低可信和故障场景用于验证节点治理，而不只是验证单一路径能跑通。",
    ]))

    slide = prs.slides.add_slide(blank)
    add_bg(slide, "方法细节：节点治理与共识提交", "4 节点、f=1、quorum=3，验证低可信排除与故障容忍")
    add_table(slide, 0.75, 1.25, 12.0, 3.0,
              ["场景", "注册节点", "可信节点", "排除节点", "不可用节点", "投票节点", "Quorum"],
              [
                  ["默认", "4", "4", "0", "0", "4", "3"],
                  ["低可信", "4", "3", "1", "0", "3", "3"],
                  ["故障", "4", "4", "0", "1", "3", "3"],
              ],
              widths=[2.0, 1.6, 1.6, 1.6, 1.8, 1.6, 1.4])
    text_box(slide, 1.0, 4.85, 11.5, 0.8, "低可信场景验证“注册不等于参与共识”；故障场景验证“可信不等于可用”。两者都要求剩余投票节点仍满足 quorum=3。", 17, WHITE, True)
    add_footer(slide, 7)
    notes.append(notes_section(7, "节点治理与共识", [
        "这里给出实验中的节点治理逻辑。默认场景 4 个节点都可信并参与投票。",
        "低可信场景中 node4 仍然注册到 ZooKeeper，但可信分低于阈值，因此被排除在共识集合之外。",
        "故障场景中 node4 注册且可信，但在共识阶段标记为不可用，剩余 3 个节点仍然可以满足 quorum。",
    ]))

    slide = prs.slides.add_slide(blank)
    add_bg(slide, "方法细节：纠删码辅助审计恢复模型", "当前是机制模型，用于验证审计发现、恢复和链上记录能否闭环")
    add_table(slide, 0.65, 1.2, 12.1, 3.2,
              ["场景", "检测率", "恢复成功率", "闭环率", "存储冗余", "链上记录"],
              [[item["label"], str(item["detection_rate"]), str(item["recovery_success_rate"]), str(item["closed_loop_rate"]), str(item["avg_storage_overhead"]), str(item["total_chain_records"])] for item in sim["recovery_summary"]],
              widths=[4.0, 1.4, 1.7, 1.4, 1.5, 1.4])
    card(slide, 1.0, 5.05, 5.4, 1.0, "关键意义", "从“发现异常”继续走到“恢复分片、链上记录、信誉惩罚”，形成审计闭环。", CYAN)
    card(slide, 7.0, 5.05, 5.4, 1.0, "当前边界", "尚未接入真实 Reed-Solomon 库；本轮先验证机制链路和指标口径。", WARN)
    add_footer(slide, 8)
    notes.append(notes_section(8, "纠删码审计恢复", [
        "这页介绍未来小论文 idea 中比较重要的一部分：纠删码辅助审计恢复。",
        "目前实现的是机制模型，不是完整 Reed-Solomon 编码库，所以讲的时候要控制预期。",
        "它的价值在于把异常发现、恢复、链上记录和信誉惩罚连接起来，这比只做哈希检查更接近数据流通的治理闭环。",
    ]))

    slide = prs.slides.add_slide(blank)
    add_bg(slide, "实验设计：真实 MNIST + 容器化联盟链", "使用标准 MNIST 60000/10000，容器内 ZooKeeper 维护节点注册")
    metric(slide, 0.75, 1.35, "训练样本", "60,000", "MNIST train", CYAN)
    metric(slide, 3.75, 1.35, "测试样本", "10,000", "MNIST test", BLUE)
    metric(slide, 6.75, 1.35, "参与方", "10", "随机切分子集", PURPLE)
    metric(slide, 9.75, 1.35, "训练轮次", "30", "第17轮异常", WARN)
    add_table(slide, 0.8, 3.15, 11.8, 2.45,
              ["组件", "实现", "作用"],
              [
                  ["ZooKeeper", "zookeeper:3.9", "节点注册和成员发现"],
                  ["链节点", "4 个 Flask 容器节点", "健康检查、投票、提交区块"],
                  ["Runner", "consortium_runner.py", "执行联邦更新并发起区块提交"],
                  ["Verifier", "verify_results.py", "校验模式、quorum、哈希链和异常轮次"],
              ],
              widths=[2.4, 3.5, 5.9])
    add_footer(slide, 9)
    notes.append(notes_section(9, "实验设计", [
        "本次真实实验使用标准 MNIST 数据，训练集 60000、测试集 10000，随机平分给 10 个参与方。",
        "容器侧包括 ZooKeeper、4 个链节点和 runner。ZooKeeper 维护节点注册表，runner 从注册表读取健康节点后发起投票。",
        "第 17 轮设置恶意反向梯度，用来验证异常更新能否被链上记录。",
    ]))

    slide = prs.slides.add_slide(blank)
    add_bg(slide, "实验流程：从数据转换到结果验证", "真实 MNIST 已转换为 data/mnist.npz，并在容器中挂载只读使用")
    flow(slide, ["IDX gzip", "mnist.npz", "Docker挂载", "FL更新", "PBFT提交", "Verifier校验"], 0.55, 1.5, 1.65, 0.28)
    add_table(slide, 0.9, 3.1, 11.4, 2.35,
              ["阶段", "命令/文件", "输出"],
              [
                  ["数据准备", "scripts/prepare_mnist_npz.py", "data/mnist.npz"],
                  ["默认容器", "docker-compose.mnist.yml", "container_results_mnist"],
                  ["低可信容器", "docker-compose.mnist.lowtrust.yml", "container_results_mnist_lowtrust"],
                  ["故障容器", "docker-compose.mnist.fault.yml", "container_results_mnist_fault"],
              ],
              widths=[2.2, 4.8, 4.4])
    add_footer(slide, 10)
    notes.append(notes_section(10, "实验流程", [
        "这页说明实验怎么跑。你提供的是 MNIST 标准 gzip IDX 文件，我先转换成项目已有加载逻辑支持的 npz 格式。",
        "为了避免重新构建镜像时访问 Docker Hub 超时，容器运行时把本地 config、src 和 data 目录挂载进去。",
        "每个场景单独输出目录，避免覆盖原来的 mnist_like 结果。",
    ]))

    slide = prs.slides.add_slide(blank)
    add_bg(slide, "初步结果 1：真实 MNIST 三组容器实验均通过", "三组实验均为 zookeeper_container，30 个区块，异常轮次=17")
    add_table(slide, 0.55, 1.25, 12.25, 2.65,
              ["场景", "模式", "区块", "可信/投票", "异常轮次", "最终准确率", "最终哈希"],
              [
                  ["默认", default["execution_mode"], "30", "4 / 4", "17", f"{default['final_accuracy']:.4f}", default["final_block_hash"][:16]],
                  ["低可信", low["execution_mode"], "30", "3 / 3", "17", f"{low['final_accuracy']:.4f}", low["final_block_hash"][:16]],
                  ["故障", fault["execution_mode"], "30", "4 / 3", "17", f"{fault['final_accuracy']:.4f}", fault["final_block_hash"][:16]],
              ],
              widths=[1.4, 2.35, 1.0, 1.4, 1.3, 1.5, 2.65])
    chart_line(slide, 0.95, 4.45, 5.5, 1.65, acc_labels, acc_curve, "默认场景准确率变化")
    chart_bar(slide, 7.05, 4.45, 5.15, 1.65, ["默认", "低可信", "故障"], [("平均提交延迟(ms)", lat_avg)], "平均提交延迟")
    add_footer(slide, 11)
    notes.append(notes_section(11, "真实 MNIST 结果", [
        "三组真实 MNIST 容器实验都通过了 verifier，说明不是单纯脚本跑完，而是哈希链、quorum、节点状态和异常轮次都被结构化校验过。",
        "最终准确率都是 0.6747，这是因为三组使用同一组模型更新和随机种子，节点治理影响的是链上确认路径，不改变模型训练路径。",
        "平均提交延迟在 20 到 23 毫秒左右，故障场景略高，但仍满足 3 节点 quorum 提交。",
    ]))

    slide = prs.slides.add_slide(blank)
    add_bg(slide, "初步结果 2：低可信与故障场景可维持提交", "目标不是提高 MNIST 精度，而是验证审计链路在节点治理场景下不断链")
    chart_bar(slide, 0.85, 1.25, 5.4, 2.35, ["默认", "低可信", "故障"], [
        ("可信节点", [4, 3, 4]),
        ("投票节点", [4, 3, 3]),
    ], "节点集合变化")
    chart_bar(slide, 7.0, 1.25, 5.3, 2.35, ["默认", "低可信", "故障"], [
        ("最大提交延迟(ms)", lat_max),
    ], "最大提交延迟")
    card(slide, 0.9, 4.4, 3.6, 1.3, "默认", "4 个可信节点参与投票，30 个区块全部提交，最终哈希 3c9840ba604f2ac9。", CYAN)
    card(slide, 4.85, 4.4, 3.6, 1.3, "低可信", "node4 被排除，剩余 3 个可信节点满足 quorum，最终哈希 578eec2e7e45341e。", GREEN)
    card(slide, 8.8, 4.4, 3.6, 1.3, "故障", "node4 不可用，剩余 3 个投票节点继续提交，最终哈希 6cea10f75ad056fa。", WARN)
    add_footer(slide, 12)
    notes.append(notes_section(12, "治理场景结果", [
        "这页强调实验目的：不是把 MNIST 精度做到很高，而是验证链上审计和节点治理机制是否能在不同场景下保持闭环。",
        "低可信节点场景体现的是基于可信评分的准入控制；故障节点场景体现的是容错能力。",
        "三个场景都提交了 30 个区块，并且都识别出第 17 轮异常更新。",
    ]))

    slide = prs.slides.add_slide(blank)
    add_bg(slide, "初步结果 3：机制仿真显示轻量化有空间", "该页是模型级估算，用于支撑后续小论文 idea，不直接等同真实链性能")
    add_table(slide, 0.65, 1.2, 12.1, 2.75,
              ["策略", "共识节点", "链上记录", "事件摊销确认(ms)", "存储冗余"],
              [[item["label"], str(item["avg_consensus_nodes"]), str(item["total_chain_records"]), str(item["avg_amortized_event_commit_ms"]), str(item["evidence_storage_overhead"])] for item in summary],
              widths=[4.8, 1.5, 1.7, 2.4, 1.5])
    metric(slide, 1.0, 4.7, "链上记录减少", f"{record_reduction*100:.2f}%", "相对 BLoP 全可信 PBFT", CYAN)
    metric(slide, 4.0, 4.7, "摊销确认降低", f"{latency_reduction*100:.2f}%", "模型级估算", GREEN)
    metric(slide, 7.0, 4.7, "存储冗余降低", f"{storage_reduction*100:.2f}%", "3.0 -> 1.5", PURPLE)
    metric(slide, 10.0, 4.7, "恢复闭环率", "0.3893", "纠删码恢复链", WARN)
    add_footer(slide, 13)
    notes.append(notes_section(13, "机制仿真结果", [
        "这页讲机制级仿真结果，注意措辞：这里是模型级估算，不是声称真实链性能已经达到这个数值。",
        "对比 BLoP 式全可信节点 PBFT，TrustFlowAudit 通过可信委员会、分级批量存证和纠删码证据模型，显著减少链上记录和存储冗余。",
        "这些结果主要用于说明后续小论文方向有优化空间，真实论文还需要进一步补强真实编码和更严格 baseline。",
    ]))

    slide = prs.slides.add_slide(blank)
    add_bg(slide, "总结与下一步", "当前已经形成可复现实验闭环，后续重点是把机制模型推进到可发表实验")
    card(slide, 0.75, 1.25, 3.7, 1.6, "已完成", "真实 MNIST 数据接入；ZooKeeper + Docker Compose 联盟链；默认、低可信、故障三组容器实验。", GREEN)
    card(slide, 4.8, 1.25, 3.7, 1.6, "可汇报结论", "链上审计路径可以在节点排除和节点故障场景下保持 30 轮提交，不依赖 TPM/TEE。", CYAN)
    card(slide, 8.85, 1.25, 3.7, 1.6, "需控制边界", "当前是 PBFT 风格原型；纠删码为机制模型；尚未替换真实 Fabric/FISCO 或 RS 编码库。", WARN)
    add_table(slide, 0.9, 3.7, 11.6, 1.9,
              ["后续任务", "目的", "优先级"],
              [
                  ["接入真实 Reed-Solomon 编码", "把纠删码从模型推进到工程实现", "P0"],
                  ["增加 baseline 对比", "全节点 PBFT、随机委员会、可信委员会", "P0"],
                  ["动态节点实验", "入网、退网、恢复后的重新准入", "P1"],
                  ["真实联盟链框架替换", "FISCO/Fabric/Tendermint 之一", "P2"],
              ],
              widths=[4.0, 5.5, 1.5])
    add_footer(slide, 14)
    notes.append(notes_section(14, "总结与下一步", [
        "最后总结当前工作：真实 MNIST 数据已经接入，容器化联盟链实验已经跑通，低可信和故障场景也验证通过。",
        "可以汇报的重点不是模型精度，而是审计路径的连通性和节点治理下的持续提交能力。",
        "下一步建议优先做两个事情：补真实 Reed-Solomon 编码，以及补全 baseline 对比。真实联盟链框架替换可以作为后续扩展，不建议现在过早放大工程范围。",
    ]))

    output_dir.mkdir(parents=True, exist_ok=True)
    pptx_path = output_dir / "TrustFlowAudit_MNIST容器实验_20分钟汇报.pptx"
    notes_path = output_dir / "TrustFlowAudit_MNIST容器实验_逐页讲稿.md"
    prs.save(pptx_path)
    notes_path.write_text("# TrustFlowAudit MNIST 容器实验逐页讲稿\n\n" + "\n".join(notes), encoding="utf-8")
    return pptx_path, notes_path


def main() -> int:
    parser = argparse.ArgumentParser(description="Generate TrustFlowAudit MNIST experiment PPT and speaker notes.")
    parser.add_argument("--output-dir", default=str(ROOT.parent / "output" / "ppt"))
    args = parser.parse_args()
    pptx_path, notes_path = build_deck(Path(args.output_dir))
    print(f"pptx: {pptx_path.resolve()}")
    print(f"notes: {notes_path.resolve()}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
